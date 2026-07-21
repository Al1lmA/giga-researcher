import os 
from io import BytesIO
import docx
from docx import Document
from loguru import logger
from pptx import Presentation
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import math
import numpy
import re

BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
OUTPUTS_DIR = os.path.join(BASE_DIR, "outputs")

    # Константы для позиционирования текста на слайде
TEXT_BOX_LEFT = Inches(0.5)
TEXT_BOX_TOP = Inches(1)
TEXT_BOX_WIDTH = Inches(6.7)
TEXT_BOX_HEIGHT = Inches(7.2)
TEXT_FONT_SIZE = Pt(14)
TEXT_FONT_NAME = "SB Sans Display"

HEADING_RE = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*$")
ORDERED_LIST_RE = re.compile(r"^(\s*)(\d+)[\.\)]\s+(.*\S)\s*$")
BULLET_LIST_RE = re.compile(r"^(\s*)[-*•]\s+(.*\S)\s*$")
LATEX_FOOTNOTE_RE = re.compile(r"\$\{\s*\}\^\{(\d+)\}\$")
MARKDOWN_FOOTNOTE_RE = re.compile(r"\[\^(\d+)[\^]?\]")
INLINE_FORMAT_RE = re.compile(r"(\*\*\*.+?\*\*\*|\*\*.+?\*\*|\*.+?\*)")
SENTENCE_SPLIT_RE = re.compile(
    r"(?<!\b[А-ЯA-Z])"
    r"(?<!\b[А-ЯA-Z]\.)"
    r"(?<!\b[а-яa-z]{1})"
    r"(?<=[.!?])\s+"
)

def safe_filename(name: str, fallback: str = "report") -> str:  ##исправление бага имени на винде, на сервер подгружать НЕ НАДО 
        """
        Делает строку безопасной для имени файла на Windows/Linux.
        Убирает символы, запрещённые в Windows: < > : " / \ | ? *
        """
        if not name:
            return fallback

        name = re.sub(r'[<>:"/\\|?*]', '', name)
        name = re.sub(r'\s+', ' ', name).strip()
        name = name.rstrip(". ")

        return name[:150] or fallback


def normalize_report_text(text: str) -> str:
        if not text:
            return ""

        normalized = text.replace("\r\n", "\n").replace("\r", "\n")
        normalized = re.sub(r"[ \t]+\n", "\n", normalized)
        normalized = re.sub(r"\n{3,}", "\n\n", normalized)
        normalized = re.sub(r"^\s*[-*•]\s+", "- ", normalized, flags=re.MULTILINE)
        normalized = re.sub(r"^\s*(\d+)[\.\)]\s*", r"\1. ", normalized, flags=re.MULTILINE)
        return normalized.strip()


def clean_report_text(text: str, preserve_inline_markup: bool = False) -> str:
        if text is None:
            return ""

        invisible_translation = str.maketrans({
            "\ufeff": " ",
            "\ufffe": " ",
            "\u2060": " ",
            "\u200b": " ",
            "\u200c": " ",
            "\u200d": " ",
            "\ufffc": " ",
        })
        cleaned = str(text).translate(invisible_translation)
        
        cleaned = cleaned.replace("\x0b", " ").replace("\xa0", " ")
        cleaned = LATEX_FOOTNOTE_RE.sub(lambda m: f"[{m.group(1)}]", cleaned)
        cleaned = MARKDOWN_FOOTNOTE_RE.sub(lambda m: f"[{m.group(1)}]", cleaned)
        cleaned = re.sub(r"[ \t]*[-]{3,}[ \t]*", " ", cleaned)
        cleaned = re.sub(r"^\s{0,3}>\s?", "", cleaned, flags=re.MULTILINE)
        if not preserve_inline_markup:
            cleaned = re.sub(r"(\*\*\*|\*\*|\*)(.+?)\1", r"\2", cleaned)
            cleaned = re.sub(r"(__|_)(.+?)\1", r"\2", cleaned)
            cleaned = cleaned.replace("**", "").replace("__", "").replace("*", "")
        cleaned = re.sub(r"[ \t]+\n", "\n", cleaned)
        cleaned = re.sub(r"\n[ \t]+", "\n", cleaned)
        cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        cleaned = re.sub(r"\$\s*\{\s*\}\s*\$", "", cleaned)
        cleaned = re.sub(r"\$\s*\{\s*\}\s*\^\s*\$", "", cleaned)
        cleaned = re.sub(r"\$\s*\$", "", cleaned)
        return cleaned.strip()


def format_display_value(value) -> str:
        if value is None:
            return "Нет данных"

        if isinstance(value, str):
            cleaned = clean_report_text(value)
            if not cleaned or cleaned.lower() == "n/a":
                return "Нет данных"

            compact = cleaned.replace(" ", "")
            if re.fullmatch(r"-?\d+(?:\.0+)?", compact):
                return f"{int(float(compact)):,}".replace(",", " ")
            return cleaned

        if isinstance(value, (int, numpy.integer)):
            return f"{int(value):,}".replace(",", " ")

        if isinstance(value, float):
            if math.isnan(value):
                return "Нет данных"
            if value.is_integer():
                return f"{int(value):,}".replace(",", " ")
            return clean_report_text(str(value))

        cleaned = clean_report_text(str(value))
        return cleaned or "Нет данных"


def split_inline_format_runs(text: str):
        prepared = clean_report_text(text, preserve_inline_markup=True)
        if not prepared:
            return []

        parts = []
        position = 0
        for match in INLINE_FORMAT_RE.finditer(prepared):
            start, end = match.span()
            if start > position:
                plain = prepared[position:start]
                if plain:
                    parts.append((plain, False, False))

            token = match.group(0)
            if token.startswith("***") and token.endswith("***"):
                parts.append((token[3:-3], True, True))
            elif token.startswith("**") and token.endswith("**"):
                parts.append((token[2:-2], True, False))
            else:
                parts.append((token[1:-1], False, True))
            position = end

        if position < len(prepared):
            tail = prepared[position:]
            if tail:
                parts.append((tail, False, False))

        return [(fragment, bold, italic) for fragment, bold, italic in parts if fragment]


def is_short_subheading(text: str, next_block=None) -> bool:
        plain = clean_report_text(text).strip()
        if not plain or not next_block:
            return False
        if len(plain) > 60 or len(plain.split()) > 6:
            return False
        if plain.endswith((".", "!", "?", ";")):
            return False
        if "," in plain:
            return False
        return next_block[0] in {"paragraph", "list"}


def write_pptx_runs(paragraph, text: str, *, base_bold: bool = False, base_italic: bool = False):
        paragraph.text = ""
        parts = split_inline_format_runs(text)
        if not parts:
            parts = [(clean_report_text(text), False, False)]

        for fragment, bold, italic in parts:
            run = paragraph.add_run()
            run.text = fragment
            run.font.size = TEXT_FONT_SIZE
            run.font.name = TEXT_FONT_NAME
            run.font.bold = base_bold or bold
            run.font.italic = base_italic or italic


def write_docx_runs(paragraph, text: str, *, base_bold: bool = False, base_italic: bool = False):
        parts = split_inline_format_runs(text)
        if not parts:
            parts = [(clean_report_text(text), False, False)]

        for fragment, bold, italic in parts:
            run = paragraph.add_run(fragment)
            run.bold = base_bold or bold
            run.italic = base_italic or italic


def iter_text_blocks(text: str):
        prepared_lines = []
        for raw_line in normalize_report_text(clean_report_text(text, preserve_inline_markup=True)).split("\n"):
            indent = len(raw_line) - len(raw_line.lstrip(" "))
            line = raw_line.strip()
            if not line:
                continue

            heading_match = HEADING_RE.match(line)
            if heading_match:
                prepared_lines.append(("heading", heading_match.group(1).strip(), 0))
                continue

            ordered_match = ORDERED_LIST_RE.match(raw_line)
            if ordered_match:
                prepared_lines.append(("list", f"{ordered_match.group(2)}. {ordered_match.group(3).strip()}", min(indent // 2, 4)))
                continue

            bullet_match = BULLET_LIST_RE.match(raw_line)
            if bullet_match:
                prepared_lines.append(("list", f"- {bullet_match.group(2).strip()}", min(indent // 2, 4)))
                continue

            prepared_lines.append(("paragraph", line, 0))

        for index, block in enumerate(prepared_lines):
            next_block = prepared_lines[index + 1] if index + 1 < len(prepared_lines) else None
            if block[0] == "paragraph" and is_short_subheading(block[1], next_block):
                yield "subheading", block[1], block[2]
            else:
                yield block


def estimate_block_units(block, chars_per_line: int) -> float:
        block_type, text, level = block
        effective_chars = chars_per_line
        if block_type == "list":
            effective_chars = max(18, chars_per_line - level * 6)
        elif block_type == "heading":
            effective_chars = max(18, chars_per_line - 8)
        elif block_type == "subheading":
            effective_chars = max(18, chars_per_line - 4)

        plain_text = clean_report_text(text)
        segments = plain_text.split("\n") if plain_text else [""]
        visual_lines = 0
        for segment in segments:
            segment = segment.strip()
            if not segment:
                visual_lines += 1
            else:
                visual_lines += max(1, math.ceil(len(segment) / effective_chars))

        spacing = 0.55 if block_type == "heading" else 0.35 if block_type == "subheading" else 0.2
        if block_type == "list":
            spacing = 0.15
        return visual_lines + spacing


def page_units(blocks, chars_per_line: int) -> float:
        return sum(estimate_block_units(block, chars_per_line) for block in blocks)

def move_trailing_heading_to_next_page(pages):
        """
        Не оставляет heading/subheading последней строкой страницы.
        Если страница заканчивается подзаголовком, переносим его
        в начало следующей страницы.
        """
        if len(pages) < 2:
            return pages

        fixed_pages = []

        for page in pages:
            if not page:
                continue

            if fixed_pages and fixed_pages[-1] and fixed_pages[-1][-1][0] in {"heading", "subheading"}:
                trailing_heading = fixed_pages[-1].pop()

                if fixed_pages[-1]:
                    page.insert(0, trailing_heading)
                else:
                    fixed_pages.pop()
                    page.insert(0, trailing_heading)

            fixed_pages.append(page)

        return fixed_pages

def split_paragraph_block(block, available_units: float, chars_per_line: int):
        block_type, text, level = block
        if block_type not in {"paragraph", "subheading"}:
            return None, block

        if available_units < 2.0:
            return None, block

        sentences = SENTENCE_SPLIT_RE.split(text)
        sentences = [sentence.strip() for sentence in sentences if sentence.strip()]
        if len(sentences) <= 1:
            sentences = re.split(r"(?<=[;:])\s+", text)
            sentences = [sentence.strip() for sentence in sentences if sentence.strip()]

        if len(sentences) <= 1:
            words = text.split()
            if not words:
                return None, block

            pieces = []
            current_words = []
            target_chars = max(chars_per_line, int(chars_per_line * max(1.0, available_units - 0.3)))
            for word in words:
                candidate = " ".join(current_words + [word]).strip()
                if current_words and len(clean_report_text(candidate)) > target_chars:
                    pieces.append(" ".join(current_words))
                    current_words = [word]
                else:
                    current_words.append(word)
            if current_words:
                pieces.append(" ".join(current_words))
            sentences = pieces

        head_parts = []
        tail_parts = []
        for sentence in sentences:
            candidate_parts = head_parts + [sentence]
            candidate_text = " ".join(candidate_parts).strip()
            candidate_block = (block_type, candidate_text, level)
            if estimate_block_units(candidate_block, chars_per_line) <= available_units:
                head_parts = candidate_parts
            else:
                tail_parts.append(sentence)

        if not head_parts or not tail_parts:
            return None, block

        head_block = (block_type, " ".join(head_parts).strip(), level)
        tail_block = (block_type, " ".join(tail_parts).strip(), level)
        return head_block, tail_block


def paginate_text_blocks(text_blocks, *, chars_per_line: int, max_units: float, min_tail_units: float = 3.0):
        work_blocks = list(text_blocks)
        if not work_blocks:
            return []

        pages = []
        current_page = []
        current_units = 0.0
        index = 0

        while index < len(work_blocks):
            block = work_blocks[index]
            next_block = work_blocks[index + 1] if index + 1 < len(work_blocks) else None
            block_units = estimate_block_units(block, chars_per_line)
            reserved_units = 0.0

            if block[0] in {"heading", "subheading"} and next_block:
                reserved_units = min(estimate_block_units(next_block, chars_per_line), 2.6)

            available_units = max_units - current_units

            if current_page and block_units + reserved_units > available_units:
                split_head, split_tail = split_paragraph_block(block, available_units, chars_per_line)
                if split_head and split_tail:
                    current_page.append(split_head)
                    current_units += estimate_block_units(split_head, chars_per_line)
                    work_blocks[index] = split_tail
                pages.append(current_page)
                current_page = []
                current_units = 0.0
                continue

            current_page.append(block)
            current_units += block_units
            index += 1

        if current_page:
            pages.append(current_page)

        pages = move_trailing_heading_to_next_page(pages)

        if len(pages) >= 2:
            last_units = page_units(pages[-1], chars_per_line)
            previous_units = page_units(pages[-2], chars_per_line)
            if last_units <= min_tail_units and previous_units + last_units <= max_units:
                pages[-2].extend(pages[-1])
                pages.pop()
                pages = move_trailing_heading_to_next_page(pages)


        return pages

class Company():
    """This is a class for the company"""
    
    def __init__(self, inn) -> None:
        self.resources = [] #источники
        self.org_name = '' #название компании
        self.text_o_kompanii = ''
        self.holders = ''
        self.products = ''
        self.inn = inn
        self.filename=f'data/{self.inn}/{self.inn}.pdf'
        self.report = ''  
        self.card = {}
        self.table = pd.DataFrame()
        self.graph = ''
        self.team = ''
        self.feedback = ''
        self.customers = ''
        self.competitors = ''
        self.trends = ''
        self.bm = ''
        self.invest = ''
        self.infra = ''
        self.conclusion = ''
        self.summ = ''

    def combine_texts(self):
        """Соединяет текстовые атрибуты в 1 строку"""
        components = [
            self.text_o_kompanii,
            self.holders,
            self.products,
            self.team,
            self.feedback,
            self.customers,
            self.competitors,
            self.trends,
            self.bm,
            self.invest,
            self.infra
        ]
        
        # Соединяем все компоненты в один текст, разделяя их пустой строкой для удобства чтения
        combined_text = "\n\n".join(filter(None, components))
        return combined_text

    def split_text(self, text, max_length=115):
        """
        Разбивает длинную строку на части по заданной максимальной длине, сохраняя абзацы и переносы строк.
        
        Args:
            text (str): Исходная длинная строка.
            max_length (int): Максимальная длина каждой части (по умолчанию 120 символов).
            
        Returns:
            list: Список частей строки.
        """
        parts = []
        parts = text.split('\n')
        # for paragraph in text.split('\n'):
        #     # parts.extend(textwrap.wrap(paragraph, width=max_length))
        #     wrapped_paragraph = textwrap.wrap(paragraph, width=max_length, drop_whitespace=False, tabsize=1)
        #     parts.extend(wrapped_paragraph)
        return parts

    async def add_text(self, prs, title, text):
        try:
            text_blocks = list(iter_text_blocks(text))
            if not text_blocks:
                return prs

            page_blocks = paginate_text_blocks(
                text_blocks,
                chars_per_line=48,
                max_units=41.0,
                min_tail_units=8.5,
            )

            for page_index, page in enumerate(page_blocks):
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                if page_index == 0:
                    slide.shapes.title.text = title
                else:
                    slide.shapes.title.text = ""

                text_box = slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE
                text_frame.paragraphs[0].font.name = TEXT_FONT_NAME
                text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                for block_type, line, level in page:
                    p = text_frame.add_paragraph()
                    write_pptx_runs(p, line, base_bold=(block_type in {"heading", "subheading"}))
                    if block_type == "list":
                        p.level = level + 1
                    

                
        except Exception as err:
            logger.error(f"Ошибка при добавлении текста: {err}")

        return prs

    async def make_pptx(self):
        """
        Формирование документа в формате pptx, по заданному шаблону.

        Структура:

        Executive summary  
        О компании  
        Анализ владения  
        Финансовые показатели   
        Бизнес-модель  
        Инвестиции 
        Ключевые клиенты 
        Продукты  
        ИТ-Инфраструктура  
        Оценка рынка и конкурентов  
        Оценка команды  
        Заключение  
        Источники  """

        # Путь к загруженному шаблону
        template_path = 'modules/template.pptx'
        # Создание нового документа на основе шаблона
        prs = Presentation(template_path)

        # Заполнение карточки компании
        slide = prs.slides[0] 
        table = slide.shapes[1].table  

        slide.shapes[0].text = f'Карточка компании\x0b«{clean_report_text(self.org_name)}»'
        title_text_frame = slide.shapes[0].text_frame
        title_text_frame.paragraphs[0].font.name = TEXT_FONT_NAME  # Задание шрифта "SB Sans Display"
        title_text_frame.paragraphs[0].font.bold = True  # Жирный шрифт
        title_text_frame.paragraphs[0].font.size = Pt(30)  # Размер шрифта 
        title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Цвет шрифта 

        data = self.card
        for row in table.rows:
            try:
                key = row.cells[0].text
                if key == 'Год основания компании':
                    # Проверяем наличие ключа в data
                    if 'Год основания компании' in data:
                        row.cells[1].text = format_display_value(data['Год основания компании'])
                    elif 'Дата регистрации компании' in data:
                        # Меняем текст на 'Дата регистрации компании' и заполняем значение
                        row.cells[0].text = 'Дата регистрации компании'
                        row.cells[0].text_frame.paragraphs[0].font.size = Pt(11)
                        row.cells[0].text_frame.paragraphs[0].font.bold = True
                        row.cells[0].text_frame.paragraphs[0].font.name = "SB Sans Display Semibold"
                        row.cells[1].text = format_display_value(data['Дата регистрации компании'])
                else:
                        row.cells[1].text = format_display_value(data.get(str(key)))
                    
            except Exception as er:
                    logger.error(er)
                    row.cells[1].text = 'Нет данных'

            row.cells[1].text_frame.paragraphs[0].font.size = Pt(11)
            row.cells[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            row.cells[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        
        # Добавление раздела "Executive summary"
        prs = await self.add_text(title='1 Executive summary', prs=prs, text=self.summ)
        
        # Добавление раздела "О компании"
        
        prs = await self.add_text(title='2. О компании', prs=prs, text=self.text_o_kompanii)
        

        # Добавление раздела "Анализ владения"
        prs = await self.add_text(title='2.1 Анализ владения', prs=prs, text=self.holders)

        try:
            # Добавление раздела "Финансовые показатели"
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = '2.2 Финансовые показатели'

            # 'Финансовые показатели'
            try:
                df = self.table
                left_inch, top_inch, width_inch, height_inch = Inches(0.3), Inches(1), Inches(6.8), Inches(4)
                new_table = slide.shapes.add_table(df.shape[0]+1, df.shape[1], left_inch, top_inch, width_inch, height_inch).table

                # Заполнение заголовков новой таблицы
                for col_num, column in enumerate(df.columns):
                    cell = new_table.cell(0, col_num)
                    cell.text = clean_report_text(column)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(13)
                            run.font.name = TEXT_FONT_NAME
                # Заполнение данных из датафрейма в новую таблицу
                for row_num in range(df.shape[0]):
                    for col_num in range(df.shape[1]):
                        value = df.iloc[row_num, col_num]
                        cell = new_table.cell(row_num+1, col_num)
                        if isinstance(value, (int, float, numpy.int64)):
                            if math.isnan(value):
                                cell.text = 'Нет данных'
                            else:
                                cell.text = "{:,.0f} тыс. ₽".format(value)
                        else:
                            cell.text = format_display_value(value)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(11)
                                run.font.name = TEXT_FONT_NAME
                # Установка ширины столбцов
                try:
                    new_table.columns[0].width = Inches(0.6) 
                    new_table.columns[1].width = Inches(1.3)
                    new_table.columns[2].width = Inches(1.5)
                    new_table.columns[3].width = Inches(1.7)
                    
                    new_table.columns[4].width = Inches(1.7)  
                except Exception as er:
                    logger.error(er)        
            except Exception as er:
                logger.error(er)
            try:
                # Добавление рисунка на слайд из BytesIO
                img_data = BytesIO(self.graph)  # Здесь должны быть данные изображения в формате BytesIO
                img = slide.shapes.add_picture(img_data, Inches(0.5), Inches(6), width=Inches(6.8), height=Inches(4))
            except Exception as er:
                logger.error(er)
        except Exception as er:
            logger.error(er)

        # Добавление раздела "Бизнес-модель"
        prs = await self.add_text(title='2.3 Бизнес-модель', prs=prs, text=self.bm)


        # Добавление раздела "Инвестиции"
        prs = await self.add_text(title='2.4 Инвестиции', prs=prs, text=self.invest)

        # Добавление раздела "Ключевые клиенты"
        prs = await self.add_text(title='2.5 Ключевые клиенты', prs=prs, text=self.customers)

        # Добавление раздела "Продукты"
        prs = await self.add_text(title='2.6 Продукты', prs=prs, text=self.products)


        # Добавление раздела "ИТ-Инфраструктура"
        prs = await self.add_text(title='2.7 ИТ-Инфраструктура', prs=prs, text=self.infra)


        # Добавление раздела "Оценка рынка и конкурентов"
        prs = await self.add_text(title='3 Оценка рынка и конкурентов', prs=prs, text=self.competitors)

        # Добавление раздела "Оценка команды"
        prs = await self.add_text(title='4 Оценка команды', prs=prs, text=self.team)

        # Добавление раздела "Заключение"
        prs = await self.add_text(title='5 Заключение', prs=prs, text=self.conclusion)

        # Добавление раздела "Источники"
        # slide = prs.slides.add_slide(prs.slide_layouts[1])
        # slide.shapes.title.text = '12. Источники'
        # text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        # text_frame = text_box.text_frame
        # p = text_frame.add_paragraph()
        # p.text = 'Текст об источниках'

        # Сохранение созданной презентации
        file_name = safe_filename(self.org_name)  #БЕЗОПАСНОЕ ИМЯ БЕЗ КАВЫЧЕК ДЛЯ ГЕНЕРАЦИИ ОТЧЕТА НА ВИНДЕ 

        pptx_path = f"{OUTPUTS_DIR}/{file_name}.pptx"
        prs.save(pptx_path)
              
        
        return pptx_path
    


    
    async def make_doc(self):
        try:
            doc = Document()

            doc.add_heading(clean_report_text(self.org_name), 0)

            doc.add_heading('Карточка компании', 1)
            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid' 
            for o in self.card:
                row_cells = table.add_row().cells
                row_cells[0].text = clean_report_text(o)
                row_cells[1].text = format_display_value(self.card[o])

        except Exception as er:
            logger.error(er)

        try:
            doc.add_heading('О компании', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.text_o_kompanii):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)

        try:
            doc.add_heading('Финансовые показатели', 1)
            t = doc.add_table(self.table.shape[0]+1, self.table.shape[1])

            for j in range(self.table.shape[-1]):
                t.cell(0,j).text = clean_report_text(self.table.columns[j])

            for i in range(self.table.shape[0]):
                for j in range(self.table.shape[-1]):
                    t.cell(i+1,j).text = format_display_value(self.table.values[i,j])

            buffer = BytesIO(self.graph)
            doc.add_picture(buffer, width = docx.shared.Cm(17))
        except Exception as er:
            logger.error(er)
        
        try:
            doc.add_heading('Продукты и услуги', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.products):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)
        
        try:
            doc.add_heading('Клиенты', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.customers):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)
        
        try:
            doc.add_heading('Отзывы', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.feedback):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)

        try:
            doc.add_heading('Вакансии', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.team):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)

        try:
            doc.add_heading('Конкуренты', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.competitors):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)

        try:
            doc.add_heading('Тренды', 1)
            for block_type, paragraph_text, _ in iter_text_blocks(self.trends):
                paragraph = doc.add_paragraph()
                write_docx_runs(paragraph, paragraph_text, base_bold=(block_type in {"heading", "subheading"}))
        except Exception as er:
            logger.error(er)
        
        # doc.add_heading('Источники', 1)
        # doc.add_paragraph("\n".join(self.resources))
        
        try:  
            path = f'data/{self.inn}/{self.org_name}.docx'
            doc.save(path)
            logger.info(f"Отчет сохранен. Путь: {path}")
        except Exception as er:
            logger.error(er)
            doc.save(f'{self.inn}.docx')
            logger.info(f"Отчет сохранен. Путь: {self.inn}.docx")

