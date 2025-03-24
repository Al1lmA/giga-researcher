import sys
import os

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import re
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import asyncio
from datetime import datetime
import locale
locale.setlocale(locale.LC_ALL, ('ru_RU', 'UTF-8'))

    # Константы для позиционирования текста на слайде
TEXT_BOX_LEFT = Inches(0.5)
TEXT_BOX_TOP = Inches(0.5)
TEXT_BOX_WIDTH = Inches(11)
TEXT_BOX_HEIGHT = Inches(2)
TEXT_FONT_SIZE = Pt(14)
TEXT_FONT_NAME = "SB Sans Display"
SOURCES_FONT_SIZE = Pt(11)



def add_footer_date(prs, slide):
    """Добавляет текущую дату в нижний колонтитул слайда."""
    date_shape = slide.shapes.add_textbox(0, prs.slide_height - Pt(30), prs.slide_width, Pt(30))
    date_frame = date_shape.text_frame
    date_frame.text = datetime.now().strftime("%d-%m-%Y")
    #  параметры шрифта для даты
    date_frame.paragraphs[0].font.size = Pt(11)  
    date_frame.paragraphs[0].font.color.rgb = RGBColor(111, 193, 178)
    date_frame.paragraphs[0].font.bold = True  
    date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT



async def add_text(prs, title, text):
    try:
        # Создаем новый слайд разделитель
        title_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = title_slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(60)
        title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # add_footer_date(prs, title_slide)

        # Создаем текстовый блок
        text_slide_layout = prs.slide_layouts[2]
        text_slide = prs.slides.add_slide(text_slide_layout)
        text_box = text_slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE
        text_frame.paragraphs[0].font.name = TEXT_FONT_NAME
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Выравнивание текста по левому краю
        add_footer_date(prs, text_slide)

        # Разбиваем текст на строки и добавляем их в текстовый блок
        text_lines = text.split('\n') # self.split_text(text)
        
        current_slide = text_slide
        current_text_frame = text_frame
        

        pattern = r'\*\*(.*?)\*\*'

        for line in text_lines:
            if len(current_text_frame.paragraphs) >= 7:  # Максимальное количество абзацев в текстовом блоке
                total_characters = sum(len(paragraph.text) for paragraph in current_text_frame.paragraphs)

                if total_characters + len(line) > 1000: # Максимальное количество символов на слайде
                    # Если текст не помещается в текстовый блок, создаем новый слайд
                    current_slide = prs.slides.add_slide(prs.slide_layouts[2])
                    current_text_frame = current_slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT).text_frame
                    current_text_frame.word_wrap = True
                    current_text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE
                    current_text_frame.paragraphs[0].font.name = TEXT_FONT_NAME
                    current_text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    add_footer_date(prs, current_slide)

            if line.startswith('#'):
                # Добавляем заголовок
                p = current_text_frame.add_paragraph()
                p.text = line.replace('#', '').strip()
                p.font.bold = True
                p.font.size = TEXT_FONT_SIZE
                p.font.name = TEXT_FONT_NAME
            
            else:
                # Добавляем обычный текст
                p = current_text_frame.add_paragraph()
                p.font.size = TEXT_FONT_SIZE
                p.font.name = TEXT_FONT_NAME
                match = re.search(pattern, line)
                if re.match(r'^\d+\.', line) or match:
                    # Если строка соответствует регулярному выражению или является нумерованным списком, добавляем ее жирным
                    text_list = line.split(':')
                    if text_list.__len__() == 1:
                        text_list = line.split('-')
                    p.text = text_list[0].replace('*', '').strip()
                    p.level = 1
                    p.font.bold = True
                    for i, item in enumerate(text_list[1:], start=2):
                        p = current_text_frame.add_paragraph()
                        p.text = item.replace('*', '').strip()
                        p.font.size = TEXT_FONT_SIZE
                        p.font.name = TEXT_FONT_NAME
                else:
                    p.text = line.strip()			
    except Exception as err:
        logger.error(f"Ошибка при добавлении текста: {err}")
    return prs




async def make_mr_pptx(task, qna_list):
    template_path = 'modules/mr/Market_Research_2.pptx'
    # Создание нового документа на основе шаблона
    prs = Presentation(template_path)
    # Титульный слайд
    slide = prs.slides[0] 
    slide.shapes.title.text =  f'{task}'.upper()
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(70)

    for qna in qna_list:
        for question, answer in qna.items():
            prs = await add_text(title=question, prs=prs, text=answer)
    
    # Сохранение созданной презентации
    BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    pptx_path = os.path.join(BASE_DIR, "outputs", "mr", f"{task}.pptx")
    # pptx_path = f"/home/TIsAmbrosyeva/giga_researcher/outputs/mr/{task}.pptx"
    try:
        prs.save(pptx_path)
    except Exception as err:
        logger.error(f"Ошибка при сохранении: {err}")
    logger.info('save prs')
    return pptx_path