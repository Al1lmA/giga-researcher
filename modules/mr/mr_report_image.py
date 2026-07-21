import os
import base64
import re
import asyncio
import locale
from datetime import datetime

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from modules.company import clean_report_text, iter_text_blocks, paginate_text_blocks, write_pptx_runs

locale.setlocale(locale.LC_ALL, ('ru_RU', 'UTF-8'))

TEXT_BOX_LEFT = Inches(0.5)
TEXT_BOX_TOP = Inches(0.5)
TEXT_BOX_WIDTH = Inches(11)
TEXT_BOX_HEIGHT = Inches(6.35)

TEXT_WITH_MAGE_BOX_WIDTH = Inches(5.5)
IMAGE_LEFT = Inches(7)
IMAGE_TOP = Inches(1)
IMAGE_WIDTH = Inches(5.5)
IMAGE_HEIGHT = Inches(4.5)
URL_FONT_SIZE = Pt(11)

TEXT_FONT_SIZE = Pt(14)
TEXT_FONT_NAME = "SB Sans Display"
SOURCES_FONT_SIZE = Pt(11)

def add_footer_date(prs, slide):
	"""Добавляет текущую дату в нижний колонтитул слайда."""
	date_shape = slide.shapes.add_textbox(0, prs.slide_height - Pt(30), prs.slide_width, Pt(30))
	date_frame = date_shape.text_frame
	date_frame.text = datetime.now().strftime("%d-%m-%Y")
	date_frame.paragraphs[0].font.size = Pt(11)
	date_frame.paragraphs[0].font.color.rgb = RGBColor(111, 193, 178)
	date_frame.paragraphs[0].font.bold = True
	date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


async def add_text(prs, title, text, image=None, url=None):
	try:
		text_blocks = list(iter_text_blocks(text))
		if not text_blocks:
			return prs

		title_slide = prs.slides.add_slide(prs.slide_layouts[1])
		title_shape = title_slide.shapes.title
		title_shape.text = clean_report_text(title)
		title_shape.text_frame.paragraphs[0].font.size = Pt(60)
		title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

		text_slide_layout = prs.slide_layouts[2]
		text_slide = prs.slides.add_slide(text_slide_layout)
		is_first_slide = bool(image)

		if image:
			image_data = base64.b64decode(image)
			with open("temp_image.jpg", "wb") as f:
				f.write(image_data)

			text_slide.shapes.add_picture("temp_image.jpg", IMAGE_LEFT, IMAGE_TOP, IMAGE_WIDTH, IMAGE_HEIGHT)
			text_box = text_slide.shapes.add_textbox(
				TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_WITH_MAGE_BOX_WIDTH, TEXT_BOX_HEIGHT
			)
		else:
			text_box = text_slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)

		if url:
			url_box = text_slide.shapes.add_textbox(
				IMAGE_LEFT, IMAGE_TOP + IMAGE_HEIGHT + Inches(0.1), IMAGE_WIDTH, Pt(20)
			)
			url_frame = url_box.text_frame
			url_frame.word_wrap = True
			p_url = url_frame.add_paragraph()
			run_link = p_url.add_run()
			run_link.text = clean_report_text(str(url))
			run_link.font.size = URL_FONT_SIZE
			run_link.hyperlink.address = str(url)

		text_frame = text_box.text_frame
		text_frame.word_wrap = True
		text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE
		text_frame.paragraphs[0].font.name = TEXT_FONT_NAME
		text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
		add_footer_date(prs, text_slide)

		if image:
			first_page_set = paginate_text_blocks(
				text_blocks,
				chars_per_line=42,
				max_units=15,
				min_tail_units=4.6,
			)
			first_page = first_page_set[:1]
			remaining_blocks = [block for page in first_page_set[1:] for block in page]
			other_pages = paginate_text_blocks(
				remaining_blocks,
				chars_per_line=88,
				max_units=32.0,
				min_tail_units=6.0,
			) if remaining_blocks else []
			page_sets = first_page + other_pages
		else:
			page_sets = paginate_text_blocks(
				text_blocks,
				chars_per_line=88,
				max_units=32.0,
				min_tail_units=6.0,
			)

		for page_index, page in enumerate(page_sets):
			if page_index == 0:
				current_text_frame = text_frame
			else:
				current_slide = prs.slides.add_slide(prs.slide_layouts[2])
				current_text_frame = current_slide.shapes.add_textbox(
					TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT
				).text_frame
				current_text_frame.word_wrap = True
				current_text_frame.paragraphs[0].font.size = TEXT_FONT_SIZE
				current_text_frame.paragraphs[0].font.name = TEXT_FONT_NAME
				current_text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
				add_footer_date(prs, current_slide)

			for block_index, (block_type, line, level) in enumerate(page):
				if block_index == 0:
					p = current_text_frame.paragraphs[0]
				else:
					p = current_text_frame.add_paragraph()

				write_pptx_runs(p, line, base_bold=(block_type in {"heading", "subheading"}))

				if block_type == "list":
					p.level = level + 1
	except Exception as err:
		logger.error(f"Ошибка при добавлении текста: {err}")
	return prs


async def make_mr_images_pptx(task, qna_list):
	template_path = 'modules/mr/Market_Research_2.pptx'
	prs = Presentation(template_path)

	slide = prs.slides[0]
	slide.shapes.title.text = clean_report_text(task).upper()
	slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(70)

	for qna in qna_list:
		for question, answer in qna.items():
			prs = await add_text(title=question, prs=prs, text=answer[0], image=answer[1], url=answer[2])

	base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
	pptx_path = os.path.join(base_dir, "outputs", "mr", f"{task}.pptx")
	prs.save(pptx_path)

	return pptx_path
