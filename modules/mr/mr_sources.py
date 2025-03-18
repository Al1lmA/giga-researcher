from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import re
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import asyncio
from datetime import datetime
import locale
from backend.utils import convert_pptx_to_pdf
locale.setlocale(locale.LC_ALL, ('ru_RU', 'UTF-8'))


# Константы для настройки текстового блока
TEXT_BOX_LEFT = Inches(1)
TEXT_BOX_TOP = Inches(0.5)
TEXT_BOX_WIDTH = Inches(12)
TEXT_BOX_HEIGHT = Inches(2)
SOURCES_FONT_SIZE = Pt(11)
TEXT_FONT_NAME = 'Arial'

def add_sources(prs, title, sources):
	try:
		# Создаем новый слайд с заголовком
		title_slide = prs.slides.add_slide(prs.slide_layouts[1])
		title_shape = title_slide.shapes.title
		title_shape.text = title
		title_shape.text_frame.paragraphs[0].font.size = Pt(60)

		# Создаем текстовый блок для списка
		text_slide_layout = prs.slide_layouts[2]
		current_slide = prs.slides.add_slide(text_slide_layout)
		text_box = current_slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
		text_frame = text_box.text_frame
		text_frame.word_wrap = True

		total_characters = 0

		for part in sources:
			for part_name, dicts in part.items():
				# Добавляем подзаголовок для каждой части
				subheading_paragraph = text_frame.add_paragraph()
				subheading_paragraph.text = f'\n{part_name}\n'
				subheading_paragraph.font.size = Pt(14)  # Размер шрифта подзаголовка
				subheading_paragraph.font.bold = True
				subheading_paragraph.alignment = PP_ALIGN.LEFT
				
				for index, d in enumerate(dicts):
					for key, value in d.items():
						entry_text = f"{index + 1}. {key}: "

						# Проверка на максимальное количество символов на слайде
						if total_characters + len(entry_text) > 3200 or len(text_frame.paragraphs) >= 15:
							# Если текст не помещается, создаем новый слайд
							current_slide = prs.slides.add_slide(text_slide_layout)
							text_box = current_slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP, TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
							text_frame = text_box.text_frame
							text_frame.word_wrap = True

							total_characters = 0  # Сброс счетчика символов

						# Добавление текста в текущий текстовый блок
						p = text_frame.add_paragraph()
						run_text = p.add_run()
						run_text.text = entry_text
						run_text.font.size = SOURCES_FONT_SIZE
						run_text.font.name = TEXT_FONT_NAME

						# Добавление ссылки как активного текста
						run_link = p.add_run()
						run_link.text = str(value)  # Преобразуем значение в строку для отображения
						run_link.font.size = SOURCES_FONT_SIZE
						run_link.hyperlink.address = str(value)  # Устанавливаем адрес гиперссылки

						total_characters += len(entry_text) + len(str(value)) + 1  # Обновляем общее количество символов

	except Exception as err:
		print(f"Ошибка при добавлении текста: {err}")

	return prs

def remove_non_unique_values_from_nested_dicts(input_list):
	# Словарь для подсчета значений
	value_count = {}

	# Подсчитываем количество вхождений каждого значения в каждом вложенном словаре
	for part in input_list:
		for key, dicts in part.items():
			for d in dicts:
				for value in d.values():
					value_count[value] = value_count.get(value, 0) + 1

	# Новый список для хранения уникальных словарей
	unique_list = []

	for part in input_list:
		new_part = {}
		for key, dicts in part.items():
			unique_dicts = []
			for d in dicts:
				# Фильтруем только уникальные значения
				unique_dict = {k: v for k, v in d.items() if value_count[v] == 1}
				if unique_dict:  # Добавляем только непустые словари
					unique_dicts.append(unique_dict)
			if unique_dicts:  # Добавляем непустой список в новый словарь
				new_part[key] = unique_dicts
		if new_part:  # Добавляем только непустые части
			unique_list.append(new_part)

	return unique_list


async def make_sources_file(task, sources):
	template_path = 'modules/mr/Market_Research_2.pptx'
	# Создание нового документа на основе шаблона
	prs = Presentation(template_path)
	# Титульный слайд
	slide = prs.slides[0] 
	slide.shapes.title.text =  f'{task}'.upper()
	slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(70)
	  
	# Очищаем список источников от неуникальных значений
	unique_sources = remove_non_unique_values_from_nested_dicts(sources)

	prs = add_sources(title="Источники", prs=prs, sources=unique_sources)
	# Сохранение созданной презентации
	pptx_path = f"/home/TIsAmbrosyeva/giga_researcher/outputs/mr/sources/Источники-{task}.pptx"
	prs.save(pptx_path)
	pdf_path = await convert_pptx_to_pdf(pptx_path, "/home/TIsAmbrosyeva/giga_researcher/outputs/mr/sources")
	return pdf_path